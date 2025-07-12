import json
import os
from flask import render_template, jsonify, request, send_file
from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib
import numpy as np
import io
import base64

# Conditional import for WeasyPrint with fallback
try:
    from weasyprint import HTML
    WEASYPRINT_AVAILABLE = True
    print("WeasyPrint is available")
except ImportError as e:
    print(f"WeasyPrint not available: {e}")
    WEASYPRINT_AVAILABLE = False
    
    # Create a mock HTML class
    class HTML:
        def __init__(self, *args, **kwargs):
            pass
        
        def write_pdf(self, output_file):
            # Create a simple placeholder PDF message
            with open(output_file, 'wb') as f:
                f.write(b'%PDF-1.4\n1 0 obj\n<</Type/Catalog/Pages 2 0 R>>\nendobj\n2 0 obj\n<</Type/Pages/Kids[3 0 R]/Count 1>>\nendobj\n3 0 obj\n<</Type/Page/Parent 2 0 R/Contents 4 0 R>>\nendobj\n4 0 obj\n<</Length 44>>stream\nBT\n/F1 12 Tf\n72 720 Td\n(PDF generation not available) Tj\nET\nendstream\nendobj\nxref\n0 5\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n0000000175 00000 n \ntrailer\n<</Size 5/Root 1 0 R>>\nstartxref\n271\n%%EOF')
            return output_file

# 한글 폰트 설정 (macOS 호환)
try:
    # macOS에서 사용 가능한 한글 폰트 시도
    import matplotlib.font_manager as fm
    
    # 사용 가능한 폰트 목록에서 한글 폰트 찾기
    font_list = [f.name for f in fm.fontManager.ttflist]
    
    # 우선순위 순으로 한글 폰트 시도
    korean_fonts = [
        'AppleGothic',      # macOS 기본 한글 폰트
        'NanumGothic',      # 나눔고딕
        'Malgun Gothic',    # 맑은 고딕 (Windows)
        'Noto Sans CJK KR', # 구글 Noto 폰트
        'DejaVu Sans',      # 기본 폰트
        'sans-serif'        # 시스템 기본 폰트
    ]
    
    # 첫 번째로 사용 가능한 폰트 선택
    selected_font = 'DejaVu Sans'  # 기본값
    for font in korean_fonts:
        if font in font_list:
            selected_font = font
            break
    
    matplotlib.rcParams['font.family'] = selected_font
    matplotlib.rcParams['axes.unicode_minus'] = False
    print(f"Using font: {selected_font}")
    
except Exception as e:
    print(f"Font configuration failed: {e}")
    # 기본 폰트 사용
    matplotlib.rcParams['font.family'] = 'DejaVu Sans'
    matplotlib.rcParams['axes.unicode_minus'] = False

class PresentationDesigner:
    """
    SEO 보고서 데이터를 기반으로 시각적 프레젠테이션을 디자인하는 클래스
    """
    
    def __init__(self, report_data):
        """
        프레젠테이션 디자이너 초기화
        
        Args:
            report_data (dict): 보고서 데이터
        """
        self.report_data = report_data
        self.charts_dir = None
        
    def generate_charts(self, output_dir):
        """
        차트 생성
        
        Args:
            output_dir (str): 차트 저장 디렉토리
            
        Returns:
            dict: 생성된 차트 파일 경로
        """
        self.charts_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
        charts = {}
        
        try:
            # 점수 차트
            charts['scores'] = self._generate_score_chart()
            print("Generated scores chart")
        except Exception as e:
            print(f"Error generating scores chart: {e}")
            charts['scores'] = None
        
        try:
            # 키워드 차트
            charts['keywords'] = self._generate_keyword_chart()
            print("Generated keywords chart")
        except Exception as e:
            print(f"Error generating keywords chart: {e}")
            charts['keywords'] = None
        
        try:
            # 페이지 깊이 분포 차트
            charts['page_depth'] = self._generate_page_depth_chart()
            print("Generated page depth chart")
        except Exception as e:
            print(f"Error generating page depth chart: {e}")
            charts['page_depth'] = None
        
        try:
            # 온페이지 SEO 점수 차트
            charts['onpage_scores'] = self._generate_onpage_scores_chart()
            print("Generated onpage scores chart")
        except Exception as e:
            print(f"Error generating onpage scores chart: {e}")
            charts['onpage_scores'] = None
        
        try:
            # 기술적 SEO 카테고리 점수 차트
            charts['technical_scores'] = self._generate_technical_scores_chart()
            print("Generated technical scores chart")
        except Exception as e:
            print(f"Error generating technical scores chart: {e}")
            charts['technical_scores'] = None
        
        return charts
    
    def _generate_score_chart(self):
        """
        SEO 점수 차트 생성
        
        Returns:
            str: 차트 파일 경로
        """
        scores = self.report_data['scores']
        
        # 데이터 준비
        categories = ['전체', '기술적 SEO', '온페이지 SEO']
        values = [scores['overall'], scores['technical'], scores['onpage']]
        
        # 색상 설정 (최소한의 색상 사용)
        colors = ['#4a6fa5', '#5b8c5a', '#d98c5f']
        
        # 그림 크기 설정
        plt.figure(figsize=(8, 5))
        
        # 바 차트 생성
        bars = plt.bar(categories, values, color=colors, width=0.6)
        
        # 바 위에 값 표시
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height + 1,
                    f'{height:.1f}',
                    ha='center', va='bottom', fontsize=12)
        
        # 축 및 제목 설정
        plt.ylim(0, 105)  # 최대값을 100보다 약간 높게 설정하여 텍스트 공간 확보
        plt.ylabel('점수 (100점 만점)', fontsize=12)
        plt.title('SEO 점수 요약', fontsize=14, fontweight='bold')
        
        # 그리드 추가 (가로선만)
        plt.grid(axis='y', linestyle='--', alpha=0.7)
        
        # 여백 조정
        plt.tight_layout()
        
        # 파일로 저장
        output_file = os.path.join(self.charts_dir, 'scores_chart.png')
        plt.savefig(output_file, dpi=100, bbox_inches='tight')
        plt.close()
        
        return output_file
    
    def _generate_keyword_chart(self):
        """
        키워드 차트 생성
        
        Returns:
            str: 차트 파일 경로
        """
        # 상위 10개 키워드 가져오기
        keywords = self.report_data['keywords']['global_keywords'][:10]
        
        # 데이터 준비
        kw_names = [kw['keyword'] for kw in keywords]
        kw_counts = [kw['count'] for kw in keywords]
        
        # 그림 크기 설정
        plt.figure(figsize=(10, 6))
        
        # 수평 바 차트 생성
        bars = plt.barh(kw_names[::-1], kw_counts[::-1], color='#4a6fa5')
        
        # 바 끝에 값 표시
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 0.5, bar.get_y() + bar.get_height()/2.,
                    f'{width}',
                    ha='left', va='center', fontsize=10)
        
        # 축 및 제목 설정
        plt.xlabel('출현 횟수', fontsize=12)
        plt.title('상위 10개 키워드', fontsize=14, fontweight='bold')
        
        # 그리드 추가 (세로선만)
        plt.grid(axis='x', linestyle='--', alpha=0.7)
        
        # 여백 조정
        plt.tight_layout()
        
        # 파일로 저장
        output_file = os.path.join(self.charts_dir, 'keywords_chart.png')
        plt.savefig(output_file, dpi=100, bbox_inches='tight')
        plt.close()
        
        return output_file
    
    def _generate_page_depth_chart(self):
        """
        페이지 깊이 분포 차트 생성
        
        Returns:
            str: 차트 파일 경로
        """
        # 사이트 구조 데이터 가져오기
        site_structure = self.report_data['technical_seo'].get('site_structure', {})
        depth_distribution = site_structure.get('depth_distribution', {})
        
        if not depth_distribution:
            # 예시 데이터 (실제 데이터가 없는 경우)
            depth_distribution = {'0': 1, '1': 5, '2': 10, '3': 7, '4': 3}
        
        # 데이터 준비
        depths = [int(d) for d in depth_distribution.keys()]
        counts = list(depth_distribution.values())
        
        # 데이터 정렬
        sorted_data = sorted(zip(depths, counts))
        depths = [d[0] for d in sorted_data]
        counts = [d[1] for d in sorted_data]
        
        # 그림 크기 설정
        plt.figure(figsize=(8, 5))
        
        # 선 그래프 생성
        plt.plot(depths, counts, marker='o', linestyle='-', color='#4a6fa5', linewidth=2, markersize=8)
        
        # 점 위에 값 표시
        for i, count in enumerate(counts):
            plt.text(depths[i], count + 0.5, str(count), ha='center', va='bottom', fontsize=10)
        
        # 축 및 제목 설정
        plt.xlabel('페이지 깊이', fontsize=12)
        plt.ylabel('페이지 수', fontsize=12)
        plt.title('페이지 깊이별 분포', fontsize=14, fontweight='bold')
        
        # x축 정수 값만 표시
        plt.xticks(depths)
        
        # 그리드 추가
        plt.grid(linestyle='--', alpha=0.7)
        
        # 여백 조정
        plt.tight_layout()
        
        # 파일로 저장
        output_file = os.path.join(self.charts_dir, 'page_depth_chart.png')
        plt.savefig(output_file, dpi=100, bbox_inches='tight')
        plt.close()
        
        return output_file
    
    def _generate_onpage_scores_chart(self):
        """
        온페이지 SEO 점수 차트 생성
        
        Returns:
            str: 차트 파일 경로
        """
        # 상위 5개 페이지 가져오기
        onpage_results = self.report_data['onpage_seo'][:5]
        
        if not onpage_results:
            # 예시 데이터 (실제 데이터가 없는 경우)
            onpage_results = [
                {'url': 'https://example.com/', 'score': 85},
                {'url': 'https://example.com/about', 'score': 78},
                {'url': 'https://example.com/products', 'score': 72},
                {'url': 'https://example.com/services', 'score': 68},
                {'url': 'https://example.com/contact', 'score': 65}
            ]
        
        # 데이터 준비
        urls = [self._shorten_url(page['url']) for page in onpage_results]
        scores = [page['score'] for page in onpage_results]
        
        # 그림 크기 설정
        plt.figure(figsize=(10, 6))
        
        # 수평 바 차트 생성
        bars = plt.barh(urls[::-1], scores[::-1], color='#5b8c5a')
        
        # 바 끝에 값 표시
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 1, bar.get_y() + bar.get_height()/2.,
                    f'{width}',
                    ha='left', va='center', fontsize=10)
        
        # 축 및 제목 설정
        plt.xlabel('점수 (100점 만점)', fontsize=12)
        plt.title('상위 5개 페이지 온페이지 SEO 점수', fontsize=14, fontweight='bold')
        
        # x축 범위 설정
        plt.xlim(0, 105)
        
        # 그리드 추가 (세로선만)
        plt.grid(axis='x', linestyle='--', alpha=0.7)
        
        # 여백 조정
        plt.tight_layout()
        
        # 파일로 저장
        output_file = os.path.join(self.charts_dir, 'onpage_scores_chart.png')
        plt.savefig(output_file, dpi=100, bbox_inches='tight')
        plt.close()
        
        return output_file
    
    def _generate_technical_scores_chart(self):
        """
        기술적 SEO 카테고리 점수 차트 생성
        
        Returns:
            str: 차트 파일 경로
        """
        # 기술적 SEO 카테고리 및 점수 (예시 데이터)
        categories = [
            'robots.txt',
            'sitemap.xml',
            '사이트 구조',
            'Core Web Vitals',
            '리다이렉트',
            '표준 링크',
            '메타 태그',
            '구조화된 데이터',
            '링크',
            '모바일 친화성',
            '보안',
            '페이지 속도'
        ]
        
        # 각 카테고리의 이슈 수를 기반으로 점수 계산 (예시)
        scores = []
        for category in categories:
            category_key = self._get_category_key(category)
            if category_key in self.report_data['technical_seo']:
                issues = self.report_data['technical_seo'][category_key].get('issues', [])
                score = max(0, 100 - len(issues) * 20)  # 이슈당 20점 감점
            else:
                score = 50  # 기본값
            scores.append(score)
        
        # 그림 크기 설정
        plt.figure(figsize=(12, 8))
        
        # 레이더 차트 준비
        angles = np.linspace(0, 2*np.pi, len(categories), endpoint=False).tolist()
        angles += angles[:1]  # 첫 번째 점을 마지막에 추가하여 폐곡선 만들기
        
        scores += scores[:1]  # 첫 번째 점을 마지막에 추가하여 폐곡선 만들기
        
        # 레이더 차트 그리기
        ax = plt.subplot(111, polar=True)
        ax.plot(angles, scores, 'o-', linewidth=2, color='#4a6fa5')
        ax.fill(angles, scores, alpha=0.25, color='#4a6fa5')
        
        # 축 설정
        ax.set_thetagrids(np.degrees(angles[:-1]), categories)
        ax.set_ylim(0, 100)
        ax.set_yticks([20, 40, 60, 80, 100])
        ax.set_yticklabels(['20', '40', '60', '80', '100'])
        ax.grid(True)
        
        # 제목 설정
        plt.title('기술적 SEO 카테고리별 점수', fontsize=14, fontweight='bold', y=1.1)
        
        # 여백 조정
        plt.tight_layout()
        
        # 파일로 저장
        output_file = os.path.join(self.charts_dir, 'technical_scores_chart.png')
        plt.savefig(output_file, dpi=100, bbox_inches='tight')
        plt.close()
        
        return output_file
    
    def _shorten_url(self, url, max_length=30):
        """
        URL 단축
        
        Args:
            url (str): 원본 URL
            max_length (int): 최대 길이
            
        Returns:
            str: 단축된 URL
        """
        if len(url) <= max_length:
            return url
            
        # 프로토콜 제거
        if '://' in url:
            url = url.split('://')[-1]
            
        # 도메인 추출
        domain = url.split('/')[0]
        
        # 경로 추출
        path = '/'.join(url.split('/')[1:])
        
        # 경로가 너무 길면 단축
        if len(domain) + len(path) + 1 > max_length:
            path_length = max_length - len(domain) - 4  # '/' 및 '...' 고려
            path = path[:path_length] + '...'
            
        return domain + '/' + path
    
    def _get_category_key(self, category):
        """
        카테고리 이름을 키로 변환
        
        Args:
            category (str): 카테고리 이름
            
        Returns:
            str: 카테고리 키
        """
        mapping = {
            'robots.txt': 'robots_txt',
            'sitemap.xml': 'sitemap',
            '사이트 구조': 'site_structure',
            'Core Web Vitals': 'core_web_vitals',
            '리다이렉트': 'redirects',
            '표준 링크': 'canonical',
            '메타 태그': 'meta_tags',
            '구조화된 데이터': 'structured_data',
            '링크': 'links',
            '모바일 친화성': 'mobile_friendly',
            '보안': 'security',
            '페이지 속도': 'page_speed'
        }
        return mapping.get(category, category.lower().replace(' ', '_'))
    
    def generate_presentation_html(self, charts, output_file):
        """
        HTML 프레젠테이션 생성
        
        Args:
            charts (dict): 차트 파일 경로
            output_file (str): 출력 파일 경로
            
        Returns:
            str: 생성된 파일 경로
        """
        # 프레젠테이션 데이터
        presentation_data = self.report_data
        
        # HTML 템플릿 생성
        html_content = self._generate_presentation_html_template(presentation_data, charts)
        
        # 파일로 저장
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(html_content)
            
        return output_file
    
    def _generate_presentation_html_template(self, data, charts):
        """
        HTML 프레젠테이션 템플릿 생성
        
        Args:
            data (dict): 프레젠테이션 데이터
            charts (dict): 차트 파일 경로
            
        Returns:
            str: HTML 템플릿
        """
        # 이미지를 base64로 인코딩
        chart_images = {}
        for key, path in charts.items():
            if path and os.path.exists(path):
                try:
                    with open(path, 'rb') as f:
                        image_data = f.read()
                        base64_data = base64.b64encode(image_data).decode('utf-8')
                        chart_images[key] = f"data:image/png;base64,{base64_data}"
                except Exception as e:
                    print(f"Error encoding chart {key}: {e}")
                    chart_images[key] = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="  # 1x1 transparent PNG
            else:
                print(f"Chart {key} not found, using placeholder")
                chart_images[key] = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="  # 1x1 transparent PNG
        
        # HTML 템플릿
        html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{data.get('title', data.get('website', {}).get('url', 'SEO Audit Report'))}</title>
    <style>
        :root {{
            --primary-color: #4a6fa5;
            --secondary-color: #5b8c5a;
            --accent-color: #d98c5f;
            --text-color: #333;
            --background-color: #fff;
            --slide-background: #f9f9f9;
        }}
        
        body {{
            font-family: 'Noto Sans KR', Arial, sans-serif;
            line-height: 1.6;
            color: var(--text-color);
            background-color: var(--background-color);
            margin: 0;
            padding: 0;
        }}
        
        .presentation-container {{
            max-width: 1000px;
            margin: 0 auto;
            padding: 20px;
        }}
        
        .controls {{
            position: fixed;
            bottom: 20px;
            left: 0;
            right: 0;
            text-align: center;
            z-index: 100;
        }}
        
        .controls button {{
            background-color: var(--primary-color);
            color: white;
            border: none;
            padding: 10px 20px;
            margin: 0 5px;
            border-radius: 5px;
            cursor: pointer;
            font-size: 16px;
        }}
        
        .controls button:hover {{
            background-color: #3a5a8c;
        }}
        
        .slide {{
            display: none;
            background-color: var(--slide-background);
            border-radius: 10px;
            padding: 40px;
            margin-bottom: 20px;
            box-shadow: 0 5px 15px rgba(0,0,0,0.1);
            min-height: 500px;
            position: relative;
        }}
        
        .slide.active {{
            display: block;
        }}
        
        .slide-number {{
            position: absolute;
            bottom: 10px;
            right: 20px;
            font-size: 14px;
            color: #999;
        }}
        
        h1, h2, h3 {{
            color: var(--primary-color);
        }}
        
        h1 {{
            font-size: 2.5em;
            margin-bottom: 20px;
            border-bottom: 2px solid var(--primary-color);
            padding-bottom: 10px;
        }}
        
        h2 {{
            font-size: 2em;
            margin-top: 0;
        }}
        
        ul, ol {{
            margin-bottom: 20px;
        }}
        
        li {{
            margin-bottom: 10px;
        }}
        
        .chart-container {{
            text-align: center;
            margin: 20px 0;
        }}
        
        .chart-container img {{
            max-width: 100%;
            height: auto;
            border-radius: 5px;
            box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        }}
        
        .score-container {{
            display: flex;
            justify-content: space-around;
            margin: 30px 0;
        }}
        
        .score-box {{
            text-align: center;
            padding: 20px;
            border-radius: 10px;
            width: 25%;
            color: white;
        }}
        
        .score-box.overall {{
            background-color: var(--primary-color);
        }}
        
        .score-box.technical {{
            background-color: var(--secondary-color);
        }}
        
        .score-box.onpage {{
            background-color: var(--accent-color);
        }}
        
        .score-value {{
            font-size: 3em;
            font-weight: bold;
            margin: 10px 0;
        }}
        
        .issues-container, .recommendations-container {{
            background-color: white;
            padding: 20px;
            border-radius: 5px;
            margin: 20px 0;
            box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        }}
        
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }}
        
        th, td {{
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }}
        
        th {{
            background-color: var(--primary-color);
            color: white;
        }}
        
        tr:nth-child(even) {{
            background-color: #f2f2f2;
        }}
        
        .download-container {{
            margin-top: 30px;
            text-align: center;
        }}
        
        .download-button {{
            display: inline-block;
            background-color: var(--primary-color);
            color: white;
            padding: 10px 20px;
            border-radius: 5px;
            text-decoration: none;
            margin: 0 10px;
        }}
        
        .download-button:hover {{
            background-color: #3a5a8c;
        }}
        
        @media print {{
            .controls, .download-container {{
                display: none;
            }}
            
            .slide {{
                display: block;
                break-after: page;
                box-shadow: none;
                min-height: auto;
            }}
        }}
    </style>
</head>
<body>
    <div class="presentation-container">
        <!-- 슬라이드 1: 제목 -->
        <div class="slide active" id="slide-1">
            <h1>{data.get('title', f"SEO Audit Report - {data.get('website', {}).get('url', 'Website')}")}</h1>
            <p><strong>분석 날짜:</strong> {data['date']}</p>
            <p><strong>웹사이트:</strong> {data['website']['url']}</p>
            <div class="chart-container">
                <img src="{chart_images['scores']}" alt="SEO 점수 차트">
            </div>
            <div class="slide-number">1 / 10</div>
        </div>
        
        <!-- 슬라이드 2: 개요 -->
        <div class="slide" id="slide-2">
            <h2>개요</h2>
            <p>이 보고서는 {data['website']['url']}의 SEO 상태에 대한 종합적인 분석을 제공합니다.</p>
            <ul>
                <li>검색 엔진 최적화(SEO)는 웹사이트의 가시성과 검색 엔진 순위를 향상시키는 과정입니다.</li>
                <li>이 감사는 기술적 SEO, 온페이지 SEO, 키워드 분석을 포함합니다.</li>
                <li>분석 날짜: {data['date']}</li>
            </ul>
            <div class="score-container">
                <div class="score-box overall">
                    <h3>전체 점수</h3>
                    <div class="score-value">{data['scores']['overall']}</div>
                    <div>/ 100</div>
                </div>
                <div class="score-box technical">
                    <h3>기술적 SEO</h3>
                    <div class="score-value">{data['scores']['technical']}</div>
                    <div>/ 100</div>
                </div>
                <div class="score-box onpage">
                    <h3>온페이지 SEO</h3>
                    <div class="score-value">{data['scores']['onpage']}</div>
                    <div>/ 100</div>
                </div>
            </div>
            <div class="slide-number">2 / 10</div>
        </div>
        
        <!-- 슬라이드 3: 주요 이슈 -->
        <div class="slide" id="slide-3">
            <h2>주요 이슈</h2>
            <div class="issues-container">
                <ul>
"""
        
        # 주요 이슈 추가
        for issue in data['summary']['top_issues'][:5]:
            html += f"                    <li>{issue}</li>\n"
            
        html += """                </ul>
            </div>
            <div class="chart-container">
                <img src="{}" alt="기술적 SEO 카테고리별 점수 차트">
            </div>
            <div class="slide-number">3 / 10</div>
        </div>
        """.format(chart_images['technical_scores'])
        
        # 슬라이드 4: 개선 권장사항
        html += """
        <!-- 슬라이드 4: 개선 권장사항 -->
        <div class="slide" id="slide-4">
            <h2>개선 권장사항</h2>
            <div class="recommendations-container">
                <ul>
"""
        
        # 개선 권장사항 추가
        for rec in data['summary']['top_recommendations'][:5]:
            html += f"                    <li>{rec}</li>\n"
            
        html += """                </ul>
            </div>
            <div class="slide-number">4 / 10</div>
        </div>
        """
        
        # 슬라이드 5: 기술적 SEO 분석
        html += """
        <!-- 슬라이드 5: 기술적 SEO 분석 -->
        <div class="slide" id="slide-5">
            <h2>기술적 SEO 분석</h2>
            <p>기술적 SEO는 검색 엔진이 웹사이트를 크롤링하고 색인화하는 방식에 영향을 미치는 요소입니다.</p>
            <table>
                <tr>
                    <th>카테고리</th>
                    <th>상태</th>
                </tr>
"""
        
        # 기술적 SEO 카테고리 상태 추가
        categories = [
            {'name': 'robots.txt', 'key': 'robots_txt', 'status_key': 'exists', 'true_text': '존재함', 'false_text': '존재하지 않음'},
            {'name': 'sitemap.xml', 'key': 'sitemap', 'status_key': 'exists', 'true_text': '존재함', 'false_text': '존재하지 않음'},
            {'name': '사이트 구조', 'key': 'site_structure', 'status_key': 'max_depth', 'prefix': '최대 깊이: '},
            {'name': 'Core Web Vitals', 'key': 'core_web_vitals', 'status_key': ['LCP', 'rating']},
            {'name': '모바일 친화성', 'key': 'mobile_friendly', 'status_key': 'is_mobile_friendly', 'true_text': '좋음', 'false_text': '개선 필요'},
            {'name': 'HTTPS', 'key': 'security', 'status_key': 'is_https', 'true_text': '사용 중', 'false_text': '사용하지 않음'}
        ]
        
        for category in categories:
            status = '정보 없음'
            
            if category['key'] in data['technical_seo']:
                if isinstance(category['status_key'], list):
                    # 중첩된 키 처리
                    value = data['technical_seo'][category['key']]
                    for key in category['status_key']:
                        if key in value:
                            value = value[key]
                        else:
                            value = None
                            break
                    status = str(value) if value is not None else '정보 없음'
                elif category['status_key'] in data['technical_seo'][category['key']]:
                    value = data['technical_seo'][category['key']][category['status_key']]
                    if isinstance(value, bool):
                        status = category['true_text'] if value else category['false_text']
                    else:
                        status = category.get('prefix', '') + str(value)
            
            html += f"""                <tr>
                    <td>{category['name']}</td>
                    <td>{status}</td>
                </tr>
"""
        
        html += """            </table>
            <div class="chart-container">
                <img src="{}" alt="페이지 깊이별 분포 차트">
            </div>
            <div class="slide-number">5 / 10</div>
        </div>
        """.format(chart_images['page_depth'])
        
        # 슬라이드 6: 상위 키워드 분석
        html += """
        <!-- 슬라이드 6: 상위 키워드 분석 -->
        <div class="slide" id="slide-6">
            <h2>상위 키워드 분석</h2>
            <p>키워드 분석은 웹사이트의 콘텐츠가 사용자의 검색 의도와 얼마나 잘 일치하는지 보여줍니다.</p>
            <div class="chart-container">
                <img src="{}" alt="상위 키워드 차트">
            </div>
            <div class="slide-number">6 / 10</div>
        </div>
        """.format(chart_images['keywords'])
        
        # 슬라이드 7: 상위 페이지
        html += """
        <!-- 슬라이드 7: 상위 페이지 -->
        <div class="slide" id="slide-7">
            <h2>상위 페이지</h2>
            <p>다음은 중요도에 따라 순위가 매겨진 상위 페이지입니다.</p>
            <table>
                <tr>
                    <th>순위</th>
                    <th>URL</th>
                    <th>점수</th>
                    <th>깊이</th>
                    <th>내부 링크 수</th>
                </tr>
"""
        
        # 상위 페이지 추가
        for i, page in enumerate(data['ranked_pages'][:5]):
            html += f"""                <tr>
                    <td>{i+1}</td>
                    <td>{self._shorten_url(page['url'])}</td>
                    <td>{page['score']}</td>
                    <td>{page['depth']}</td>
                    <td>{page['inbound_links']}</td>
                </tr>
"""
        
        html += """            </table>
            <div class="slide-number">7 / 10</div>
        </div>
        """
        
        # 슬라이드 8: 온페이지 SEO 분석
        html += """
        <!-- 슬라이드 8: 온페이지 SEO 분석 -->
        <div class="slide" id="slide-8">
            <h2>온페이지 SEO 분석</h2>
            <p>온페이지 SEO는 개별 페이지의 콘텐츠와 HTML 소스 코드를 최적화하는 것을 의미합니다.</p>
            <div class="chart-container">
                <img src="{}" alt="온페이지 SEO 점수 차트">
            </div>
""".format(chart_images['onpage_scores'])
        
        # 상위 3개 페이지의 주요 이슈 추가
        if data['onpage_seo']:
            html += """            <h3>주요 페이지 이슈</h3>
            <table>
                <tr>
                    <th>페이지</th>
                    <th>주요 이슈</th>
                </tr>
"""
            
            for page in data['onpage_seo'][:3]:
                issues = ', '.join(page['issues'][:2]) if page['issues'] else '이슈 없음'
                html += f"""                <tr>
                    <td>{self._shorten_url(page['url'])}</td>
                    <td>{issues}</td>
                </tr>
"""
            
            html += """            </table>
"""
        
        html += """            <div class="slide-number">8 / 10</div>
        </div>
        """
        
        # 슬라이드 9: 다음 단계
        html += """
        <!-- 슬라이드 9: 다음 단계 -->
        <div class="slide" id="slide-9">
            <h2>다음 단계</h2>
            <p>다음은 SEO를 개선하기 위한 권장 단계입니다.</p>
            
            <h3>우선순위가 높은 작업</h3>
            <ol>
"""
        
        # 우선순위가 높은 작업 추가
        for rec in data['summary']['top_recommendations'][:3]:
            html += f"                <li>{rec}</li>\n"
            
        html += """            </ol>
            
            <h3>중기 작업</h3>
            <ol>
"""
        
        # 중기 작업 추가
        for rec in data['summary']['top_recommendations'][3:6]:
            html += f"                <li>{rec}</li>\n"
            
        html += """            </ol>
            
            <h3>장기 작업</h3>
            <ol>
"""
        
        # 장기 작업 추가
        for rec in data['summary']['top_recommendations'][6:9]:
            html += f"                <li>{rec}</li>\n"
            
        html += """            </ol>
            <div class="slide-number">9 / 10</div>
        </div>
        """
        
        # 슬라이드 10: 감사합니다
        html += """
        <!-- 슬라이드 10: 감사합니다 -->
        <div class="slide" id="slide-10">
            <h2>감사합니다</h2>
            <p>이 SEO 감사 보고서가 웹사이트의 검색 엔진 최적화를 개선하는 데 도움이 되기를 바랍니다.</p>
            <p>질문이 있으시면 언제든지 문의해 주세요.</p>
            
            <div class="download-container">
                <a href="/download/pptx" class="download-button">PPTX 다운로드</a>
                <a href="/download/pdf" class="download-button">PDF 다운로드</a>
            </div>
            <div class="slide-number">10 / 10</div>
        </div>
    </div>
    
    <div class="controls">
        <button id="prev-btn">이전</button>
        <button id="next-btn">다음</button>
    </div>
    
    <script>
        // 슬라이드 제어
        const slides = document.querySelectorAll('.slide');
        let currentSlide = 0;
        
        function showSlide(index) {
            slides.forEach(slide => slide.classList.remove('active'));
            slides[index].classList.add('active');
            currentSlide = index;
            
            // 버튼 상태 업데이트
            document.getElementById('prev-btn').disabled = currentSlide === 0;
            document.getElementById('next-btn').disabled = currentSlide === slides.length - 1;
        }
        
        document.getElementById('prev-btn').addEventListener('click', () => {
            if (currentSlide > 0) {
                showSlide(currentSlide - 1);
            }
        });
        
        document.getElementById('next-btn').addEventListener('click', () => {
            if (currentSlide < slides.length - 1) {
                showSlide(currentSlide + 1);
            }
        });
        
        // 키보드 제어
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowLeft') {
                if (currentSlide > 0) {
                    showSlide(currentSlide - 1);
                }
            } else if (e.key === 'ArrowRight') {
                if (currentSlide < slides.length - 1) {
                    showSlide(currentSlide + 1);
                }
            }
        });
        
        // 초기 상태 설정
        showSlide(0);
    </script>
</body>
</html>
"""
        
        return html
    
    def generate_pptx(self, charts, output_file):
        """
        PPTX 프레젠테이션 생성
        
        Args:
            charts (dict): 차트 파일 경로
            output_file (str): 출력 파일 경로
            
        Returns:
            str: 생성된 파일 경로
        """
        # 프레젠테이션 생성
        prs = Presentation()
        
        # 슬라이드 추가
        self._add_title_slide(prs)
        self._add_overview_slide(prs)
        self._add_issues_slide(prs)
        self._add_recommendations_slide(prs)
        self._add_technical_seo_slide(prs, charts)
        self._add_keywords_slide(prs, charts)
        self._add_top_pages_slide(prs)
        self._add_onpage_seo_slide(prs, charts)
        self._add_next_steps_slide(prs)
        self._add_thank_you_slide(prs)
        
        # 파일로 저장
        prs.save(output_file)
        
        return output_file
    
    def _add_title_slide(self, prs):
        """
        제목 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # 제목 및 부제목 설정
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self.report_data.get('title', f"SEO Audit Report - {self.report_data.get('website', {}).get('url', 'Website')}")
        subtitle.text = f"분석 날짜: {self.report_data['date']}\n웹사이트: {self.report_data['website']['url']}"
    
    def _add_overview_slide(self, prs):
        """
        개요 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "개요"
        content.text = f"""이 보고서는 {self.report_data['website']['url']}의 SEO 상태에 대한 종합적인 분석을 제공합니다.
        
• 검색 엔진 최적화(SEO)는 웹사이트의 가시성과 검색 엔진 순위를 향상시키는 과정입니다.
• 이 감사는 기술적 SEO, 온페이지 SEO, 키워드 분석을 포함합니다.
• 분석 날짜: {self.report_data['date']}

전체 점수: {self.report_data['scores']['overall']}/100
기술적 SEO 점수: {self.report_data['scores']['technical']}/100
온페이지 SEO 점수: {self.report_data['scores']['onpage']}/100"""
    
    def _add_issues_slide(self, prs):
        """
        주요 이슈 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "주요 이슈"
        
        issues_text = ""
        for issue in self.report_data['summary']['top_issues'][:5]:
            issues_text += f"• {issue}\n"
            
        content.text = issues_text
    
    def _add_recommendations_slide(self, prs):
        """
        개선 권장사항 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "개선 권장사항"
        
        recommendations_text = ""
        for rec in self.report_data['summary']['top_recommendations'][:5]:
            recommendations_text += f"• {rec}\n"
            
        content.text = recommendations_text
    
    def _add_technical_seo_slide(self, prs, charts):
        """
        기술적 SEO 분석 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
            charts (dict): 차트 파일 경로
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "기술적 SEO 분석"
        content.text = """기술적 SEO는 검색 엔진이 웹사이트를 크롤링하고 색인화하는 방식에 영향을 미치는 요소입니다.

• robots.txt: 검색 엔진 크롤러에게 웹사이트의 어떤 부분을 크롤링해야 하는지 알려주는 파일
• sitemap.xml: 웹사이트의 모든 페이지 목록을 제공하여 검색 엔진이 콘텐츠를 더 효율적으로 크롤링할 수 있도록 돕는 파일
• Core Web Vitals: 사용자 경험을 측정하는 Google의 지표 (LCP, FID, CLS)
• 모바일 친화성: 웹사이트가 모바일 기기에서 얼마나 잘 작동하는지 측정"""
        
        # 차트 추가
        if 'technical_scores' in charts and charts['technical_scores'] and os.path.exists(charts['technical_scores']):
            try:
                left = prs.slide_width * 0.1
                top = prs.slide_height * 0.5
                width = prs.slide_width * 0.8
                height = prs.slide_height * 0.4
                
                slide.shapes.add_picture(charts['technical_scores'], left, top, width, height)
            except Exception as e:
                print(f"Error adding technical scores chart to PPTX: {e}")
    
    def _add_keywords_slide(self, prs, charts):
        """
        상위 키워드 분석 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
            charts (dict): 차트 파일 경로
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "상위 키워드 분석"
        content.text = "키워드 분석은 웹사이트의 콘텐츠가 사용자의 검색 의도와 얼마나 잘 일치하는지 보여줍니다."
        
        # 차트 추가
        if 'keywords' in charts and charts['keywords'] and os.path.exists(charts['keywords']):
            try:
                left = prs.slide_width * 0.1
                top = prs.slide_height * 0.3
                width = prs.slide_width * 0.8
                height = prs.slide_height * 0.6
                
                slide.shapes.add_picture(charts['keywords'], left, top, width, height)
            except Exception as e:
                print(f"Error adding keywords chart to PPTX: {e}")
    
    def _add_top_pages_slide(self, prs):
        """
        상위 페이지 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "상위 페이지"
        
        pages_text = "다음은 중요도에 따라 순위가 매겨진 상위 페이지입니다.\n\n"
        
        for i, page in enumerate(self.report_data['ranked_pages'][:5]):
            pages_text += f"{i+1}. {self._shorten_url(page['url'])}\n"
            pages_text += f"   점수: {page['score']}, 깊이: {page['depth']}, 내부 링크 수: {page['inbound_links']}\n\n"
            
        content.text = pages_text
    
    def _add_onpage_seo_slide(self, prs, charts):
        """
        온페이지 SEO 분석 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
            charts (dict): 차트 파일 경로
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "온페이지 SEO 분석"
        content.text = "온페이지 SEO는 개별 페이지의 콘텐츠와 HTML 소스 코드를 최적화하는 것을 의미합니다."
        
        # 차트 추가
        if 'onpage_scores' in charts and charts['onpage_scores'] and os.path.exists(charts['onpage_scores']):
            try:
                left = prs.slide_width * 0.1
                top = prs.slide_height * 0.3
                width = prs.slide_width * 0.8
                height = prs.slide_height * 0.6
                
                slide.shapes.add_picture(charts['onpage_scores'], left, top, width, height)
            except Exception as e:
                print(f"Error adding onpage scores chart to PPTX: {e}")
    
    def _add_next_steps_slide(self, prs):
        """
        다음 단계 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 및 내용 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 제목 및 내용 설정
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "다음 단계"
        
        steps_text = "다음은 SEO를 개선하기 위한 권장 단계입니다.\n\n"
        
        steps_text += "우선순위가 높은 작업:\n"
        for i, rec in enumerate(self.report_data['summary']['top_recommendations'][:3]):
            steps_text += f"{i+1}. {rec}\n"
            
        steps_text += "\n중기 작업:\n"
        for i, rec in enumerate(self.report_data['summary']['top_recommendations'][3:6]):
            steps_text += f"{i+1}. {rec}\n"
            
        steps_text += "\n장기 작업:\n"
        for i, rec in enumerate(self.report_data['summary']['top_recommendations'][6:9]):
            steps_text += f"{i+1}. {rec}\n"
            
        content.text = steps_text
    
    def _add_thank_you_slide(self, prs):
        """
        감사합니다 슬라이드 추가
        
        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        # 제목 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # 제목 및 부제목 설정
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "감사합니다"
        subtitle.text = "이 SEO 감사 보고서가 웹사이트의 검색 엔진 최적화를 개선하는 데 도움이 되기를 바랍니다.\n\n질문이 있으시면 언제든지 문의해 주세요."
    
    def generate_pdf(self, html_file, output_file):
        """
        PDF 프레젠테이션 생성
        
        Args:
            html_file (str): HTML 파일 경로
            output_file (str): 출력 파일 경로
            
        Returns:
            str: 생성된 파일 경로
        """
        if WEASYPRINT_AVAILABLE:
            # HTML을 PDF로 변환
            HTML(filename=html_file).write_pdf(output_file)
        else:
            # WeasyPrint가 없으면 placeholder PDF 생성
            print("WeasyPrint not available, creating placeholder PDF")
            html_obj = HTML()
            html_obj.write_pdf(output_file)
        
        return output_file
